import io
import os
import json
import logging
from typing import List, Dict, Any
from flask import Flask, request, render_template
from pptx import Presentation
from pptx.util import Inches
import requests

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 25 * 1024 * 1024  # 25 MB

logging.getLogger("werkzeug").setLevel(logging.WARNING)
app.logger.setLevel(logging.WARNING)

ALLOWED_EXTENSIONS = {"pptx", "potx"}

def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

@app.get("/")
def index():
    return render_template("index.html")

@app.post("/generate")
def generate():
    input_text = request.form.get("input_text", "").strip()
    guidance = request.form.get("guidance", "").strip()
    provider = request.form.get("provider", "openai").strip()
    model = request.form.get("model", "").strip()
    with_notes = request.form.get("with_notes", "false").strip().lower() == "true"
    api_key = request.form.get("api_key", "").strip()

    if not input_text or not api_key or provider not in {"openai", "anthropic", "gemini"}:
        return "<div class='text-red-600'>Missing required fields.</div>", 400

    file = request.files.get("template_file")
    if file is None or file.filename == "":
        return "<div class='text-red-600'>Please upload a .pptx or .potx file.</div>", 400
    if not allowed_file(file.filename):
        return "<div class='text-red-600'>Unsupported file type.</div>", 400

    file_bytes = file.read()
    try:
        template = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        return f"<div class='text-red-600'>Failed to read PowerPoint: {e}</div>", 400

    images = extract_images_from_presentation(template)

    outline = llm_build_outline(
        provider=provider,
        model=model,
        api_key=api_key,
        input_text=input_text,
        guidance=guidance,
        with_notes=with_notes,
    )
    if "error" in outline:
        return f"<div class='text-red-600'>LLM error: {outline['error']}</div>", 400

    output_prs = Presentation(io.BytesIO(file_bytes))
    add_slides_from_outline(output_prs, outline, images)

    out = io.BytesIO()
    output_prs.save(out)
    out.seek(0)

    html = f"""
      <div class="rounded-2xl bg-green-50 border border-green-200 p-4">
        <h3 class="font-semibold text-green-700 mb-2">Success!</h3>
        <p class="text-sm text-green-800 mb-3">Your presentation is ready.</p>
        <a href="/download?name=generated.pptx" hx-boost="true" class="inline-block px-4 py-2 rounded-xl bg-green-600 text-white hover:bg-green-500">Download .pptx</a>
      </div>
      <template id="pptx-bytes">{json.dumps(list(out.getvalue()))}</template>
      <script>
        const bytes = new Uint8Array(JSON.parse(document.querySelector('#pptx-bytes').innerHTML));
        window.__PPTX_BLOB__ = new Blob([bytes], {{ type: 'application/vnd.openxmlformats-officedocument.presentationml.presentation' }});
      </script>
    """
    return html

@app.get("/download")
def download():
    return """
      <html><body>
      <script>
        if (!window.opener || !window.opener.__PPTX_BLOB__) {
          document.write('No file available. Please generate again.');
        } else {
          const a = document.createElement('a');
          a.href = URL.createObjectURL(window.opener.__PPTX_BLOB__);
          a.download = new URLSearchParams(location.search).get('name') || 'generated.pptx';
          a.click();
          setTimeout(()=>window.close(), 500);
          document.write('Downloading... You can close this tab.');
        }
      </script>
      </body></html>
    """

def extract_images_from_presentation(prs: Presentation) -> List[bytes]:
    images = []
    try:
        part = prs.part
        for rel in part.rels.values():
            target = rel._target
            try:
                if hasattr(target, "blob"):
                    images.append(target.blob)
            except Exception:
                continue
    except Exception:
        pass
    return images

def add_slides_from_outline(prs: Presentation, outline: Dict[str, Any], images: List[bytes]) -> None:
    layouts = list(prs.slide_layouts)
    if not layouts:
        return
    img_ix = 0
    for idx, slide_def in enumerate(outline.get("slides", [])):
        layout = choose_layout(layouts, idx)
        slide = prs.slides.add_slide(layout)

        title_text = slide_def.get("title", "")
        bullets = slide_def.get("bullets", [])
        notes = slide_def.get("notes", "")

        if slide.shapes.title:
            slide.shapes.title.text = title_text[:255]

        body_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:
                body_placeholder = shape
                break

        if body_placeholder and body_placeholder.has_text_frame:
            tf = body_placeholder.text_frame
            tf.clear()
            for i, b in enumerate(bullets):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = b.strip()
                p.level = 0

        if images:
            placed = False
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 18:
                    try:
                        image_blob = images[img_ix % len(images)]
                        shape.insert_picture(io.BytesIO(image_blob))
                        placed = True
                        img_ix += 1
                        break
                    except Exception:
                        continue
            if not placed:
                try:
                    image_blob = images[img_ix % len(images)]
                    left = prs.slide_width - Inches(3.0)
                    top = Inches(1.5)
                    slide.shapes.add_picture(io.BytesIO(image_blob), left, top, width=Inches(2.5))
                    img_ix += 1
                except Exception:
                    pass

        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

def choose_layout(layouts, idx: int):
    if len(layouts) == 0:
        return layouts[0]
    if idx == 0:
        for l in layouts:
            name = (l.name or "").lower()
            if "title" in name and "content" not in name:
                return l
        return layouts[0]
    for l in layouts:
        name = (l.name or "").lower()
        if "title" in name and ("content" in name or "text" in name):
            return l
    return layouts[idx % len(layouts)]

def llm_build_outline(provider: str, model: str, api_key: str, input_text: str, guidance: str, with_notes: bool) -> Dict[str, Any]:
    system = (
        "You are a slide-structuring assistant. "
        "Given a long input text (markdown or prose) and optional guidance, "
        "produce a concise slide outline as JSON with this exact schema:\n"
        "{ 'slides': [ { 'title': str, 'bullets': [str, ...], 'notes': str } , ...] }\n"
        "Aim for 6–15 slides depending on content, avoid duplication, keep bullets short (max ~12 words). "
        "If notes are not requested, set 'notes' to an empty string."
    )
    user = f"GUIDANCE: {guidance or '(none)'}\n\nTEXT:\n{input_text}"

    try:
        if provider == "openai":
            return call_openai_chat(api_key, model or "gpt-4o-mini", system, user, with_notes)
        elif provider == "anthropic":
            return call_anthropic_messages(api_key, model or "claude-3-5-sonnet-20240620", system, user, with_notes)
        elif provider == "gemini":
            return call_gemini_generate(api_key, model or "gemini-1.5-pro", system, user, with_notes)
        else:
            return {"error": "Unsupported provider"}
    except Exception as e:
        return {"error": str(e)}

def coerce_outline(text: str, with_notes: bool) -> Dict[str, Any]:
    try:
        start = text.index("{")
        end = text.rindex("}") + 1
        obj = json.loads(text[start:end])
        slides = obj.get("slides", [])
        norm = []
        for s in slides:
            norm.append({
                "title": str(s.get("title", ""))[:255],
                "bullets": [str(b)[:255] for b in (s.get("bullets") or [])][:10],
                "notes": str(s.get("notes", ""))[:2000] if with_notes else ""
            })
        if len(norm) < 1:
            norm = [{"title": "Overview", "bullets": ["Summary"], "notes": ""}]
        if len(norm) > 25:
            norm = norm[:25]
        return {"slides": norm}
    except Exception:
        return {"slides": [{"title": "Overview", "bullets": [text[:120]], "notes": "" if not with_notes else text[:500]}]}

def call_openai_chat(api_key: str, model: str, system: str, user: str, with_notes: bool) -> Dict[str, Any]:
    url = "https://api.openai.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    data = {
        "model": model,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        "temperature": 0.3,
    }
    r = requests.post(url, headers=headers, json=data, timeout=60)
    r.raise_for_status()
    content = r.json()["choices"][0]["message"]["content"]
    return coerce_outline(content, with_notes)

def call_anthropic_messages(api_key: str, model: str, system: str, user: str, with_notes: bool) -> Dict[str, Any]:
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        "x-api-key": api_key,
        "anthropic-version": "2023-06-01",
        "content-type": "application/json",
    }
    data = {
        "model": model,
        "max_tokens": 2048,
        "system": system,
        "messages": [{"role": "user", "content": user}],
        "temperature": 0.3,
    }
    r = requests.post(url, headers=headers, json=data, timeout=60)
    r.raise_for_status()
    content = r.json()["content"][0]["text"]
    return coerce_outline(content, with_notes)

def call_gemini_generate(api_key: str, model: str, system: str, user: str, with_notes: bool) -> Dict[str, Any]:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    data = {
        "contents": [
            {"role": "user", "parts": [{"text": system + "\n\n" + user}]}
        ],
        "generationConfig": {"temperature": 0.3},
    }
    r = requests.post(url, json=data, timeout=60)
    r.raise_for_status()
    content = r.json()["candidates"][0]["content"]["parts"][0]["text"]
    return coerce_outline(content, with_notes)

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8000, debug=False)
